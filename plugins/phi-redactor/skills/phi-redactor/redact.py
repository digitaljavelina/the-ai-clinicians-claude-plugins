"""
redact.py - De-identify a real document.

This is the tool you actually reach for: hand it a note, get back a redacted
copy. It reads .txt, .md, .pdf, and .docx, runs the redaction engine, and writes
a plain-text redacted file. The clinical engine (Presidio + obi/deid_roberta_i2b2)
is the default because it catches names in prose that the baseline misses.

    uv run redact.py inputs/sample_note.txt         # clinical engine (default)
    uv run redact.py inputs/sample_note.txt --engine baseline
    uv run redact.py inputs/note.pdf                # PDF -> redacted .txt
    uv run redact.py inputs/note.docx -o clean.txt  # pick the output path
    uv run redact.py inputs/                         # every supported file in a folder
    cat note.txt | uv run redact.py -               # stdin -> stdout

Safety notes:
  * Output is ALWAYS plain text. A redacted .pdf/.docx that still carried the
    original bytes in metadata or revision history would be a leak; extracting to
    text and writing text avoids that whole class of mistake.
  * The summary prints entity COUNTS only, never the detected values. This is a
    PHI tool; it must not echo PHI to your terminal or logs.
  * Over-redaction is the safe error; under-redaction is a breach. Always eyeball
    the output before trusting it. The measurement lab (evaluate.py) is how you
    justify that trust with numbers.
"""
import argparse
import json
import os
import re
import sys
from collections import Counter
from pathlib import Path

# Force the Hugging Face libraries offline for every redaction run. After bootstrap
# the models are already on disk, so redaction needs no network at all. Making a
# network call impossible here means a note being redacted can never leave the
# machine, not even via the harmless model-version check the library does at load.
# These are set BEFORE transformers is imported (engine.py imports it lazily) so
# they take effect. bootstrap.py deliberately does NOT set them: installation is
# the one time a download is allowed. setdefault lets an expert override with
# HF_HUB_OFFLINE=0 if they ever truly need online behavior.
os.environ.setdefault("HF_HUB_OFFLINE", "1")
os.environ.setdefault("TRANSFORMERS_OFFLINE", "1")
os.environ.setdefault("HF_HUB_DISABLE_TELEMETRY", "1")

# Read and write UTF-8 no matter the console default. Windows consoles default to
# a legacy code page that would mangle non-ASCII text in a note.
for _stream in (sys.stdin, sys.stdout, sys.stderr):
    try:
        _stream.reconfigure(encoding="utf-8")
    except (AttributeError, ValueError):
        pass

import engine as E  # noqa: E402  (must follow the offline env setup above)

SUPPORTED = {".txt", ".md", ".pdf", ".docx", ".xlsx", ".pptx"}

SKILL_DIR = Path(__file__).resolve().parent
MARKER = SKILL_DIR / ".installed"


# --------------------------------------------------------------------------- #
# Text extraction: turn a file of any supported type into a plain string.
# --------------------------------------------------------------------------- #
def _read_text(path: Path) -> str:
    """UTF-8 first, fall back to latin-1 so an odd export still reads."""
    data = path.read_bytes()
    try:
        return data.decode("utf-8")
    except UnicodeDecodeError:
        return data.decode("latin-1")


def _normalize_extracted(text: str) -> str:
    """Tidy text pulled from PDF/DOCX: collapse the space-padding that fixed-width
    PDF extraction adds, and drop spaces sitting just before a line break. This
    does NOT try to un-wrap words that were split across lines (guessing at that
    is unsafe and can hide PHI from the model); residual_scan() is what flags any
    identifier a line-wrap left behind."""
    text = re.sub(r"[ \t]+", " ", text)
    text = re.sub(r" *\n", "\n", text)
    return text


def _read_pdf(path: Path) -> str:
    from pypdf import PdfReader

    reader = PdfReader(str(path))
    raw = "\n".join((page.extract_text() or "") for page in reader.pages)
    return _normalize_extracted(raw)


def _read_docx(path: Path) -> str:
    from docx import Document

    doc = Document(str(path))
    parts = [p.text for p in doc.paragraphs]
    # Table cells hold PHI too (demographics grids, contact tables), so pull them.
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                parts.append(cell.text)
    return _normalize_extracted("\n".join(parts))


def _read_xlsx(path: Path) -> str:
    """Every cell across every sheet, plus the sheet tab names (a tab called
    'Smith labs' is PHI). One cell per line so distinct identifiers stay separate
    tokens. Cell comments, headers/footers, and embedded charts are NOT captured."""
    from openpyxl import load_workbook

    wb = load_workbook(filename=str(path), read_only=True, data_only=True)
    parts = []
    for ws in wb.worksheets:
        parts.append(ws.title)
        for row in ws.iter_rows(values_only=True):
            for cell in row:
                if cell is not None:
                    parts.append(str(cell))
    wb.close()
    return _normalize_extracted("\n".join(parts))


def _read_pptx(path: Path) -> str:
    """Text from every slide: text boxes, tables, grouped shapes, and speaker
    notes. Chart data and embedded objects are NOT captured."""
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    parts = []

    def walk(shapes):
        for shape in shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                walk(shape.shapes)
                continue
            if shape.has_text_frame:
                parts.append(shape.text_frame.text)
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        parts.append(cell.text)

    prs = Presentation(str(path))
    for slide in prs.slides:
        walk(slide.shapes)
        if slide.has_notes_slide:
            parts.append(slide.notes_slide.notes_text_frame.text)
    return _normalize_extracted("\n".join(parts))


def extract(path: Path) -> str:
    ext = path.suffix.lower()
    if ext == ".pdf":
        return _read_pdf(path)
    if ext == ".docx":
        return _read_docx(path)
    if ext == ".xlsx":
        return _read_xlsx(path)
    if ext == ".pptx":
        return _read_pptx(path)
    return _read_text(path)  # .txt, .md, and anything else texty


# --------------------------------------------------------------------------- #
# Extraction-sanity guard.
#
# This is the one safety decision worth owning deliberately. A scanned PDF (an
# image of a fax) or a corrupt .docx can extract to an EMPTY or near-empty
# string. If we redact "" we produce a clean-looking output file while the
# original document still carries every identifier. That silent failure is the
# most dangerous thing this tool can do, so we fail closed instead of trusting a
# suspiciously thin extraction.
#
# Tune MIN_CHARS to your real documents: a legitimate short note (a one-line
# telephone encounter) versus a scanned page that yielded nothing.
# --------------------------------------------------------------------------- #
MIN_CHARS = 20


def verify_extraction(text: str, source: Path) -> None:
    """Raise if extraction almost certainly failed to read the document."""
    stripped = text.strip()
    if not stripped:
        raise RuntimeError(
            f"{source.name}: extracted no text. If this is a scanned/image PDF, "
            f"it has NOT been redacted (the image still contains everything). "
            f"OCR it to text first, then redact."
        )
    if len(stripped) < MIN_CHARS:
        raise RuntimeError(
            f"{source.name}: extracted only {len(stripped)} characters, which "
            f"looks like a failed read rather than a real note. Refusing to write "
            f"an output that would look redacted but never saw the source. "
            f"Lower MIN_CHARS in redact.py if this document really is that short."
        )


# --------------------------------------------------------------------------- #
# Redaction + output
# --------------------------------------------------------------------------- #
def redact_text(text: str, analyzer) -> tuple[str, Counter]:
    return E.redact(text, analyzer)


# Post-redaction safety net. In correctly de-identified clinical text these
# patterns should not appear; if one does, a line-wrap or edge case likely left a
# fragment of an identifier behind, so we WARN and tell the user to eyeball it.
_RESIDUAL = {
    "email": re.compile(r"\w@|@\w"),
    "ssn": re.compile(r"\b\d{3}-\d{2}-\d{4}\b"),
    "phone": re.compile(r"\(\d{3}\)\s*\d{3}|\b\d{3}[.\-\s]\d{3}[.\-\s]\d{4}\b"),
}


def residual_scan(text: str) -> list[str]:
    return sorted(name for name, pat in _RESIDUAL.items() if pat.search(text))


def _out_path(src: Path, out_dir: Path, used: set) -> Path:
    """Name the output <stem>.redacted.txt, but never let two source files with
    the same stem (note.pdf and note.txt) collide onto one output and silently
    overwrite each other. On collision, fold the source extension into the name."""
    dest = out_dir / f"{src.stem}.redacted.txt"
    if dest in used:
        dest = out_dir / f"{src.stem}.{src.suffix.lstrip('.')}.redacted.txt"
    used.add(dest)
    return dest


def _write_map(dest: Path, src: Path, mapping: dict) -> Path:
    """Write the placeholder->original key next to the redacted file. This file is
    the re-identification key and contains PHI; it must never be sent to an AI."""
    map_path = dest.with_suffix(".map.json")
    map_path.write_text(
        json.dumps(
            {
                "source": src.name,
                "placeholder_format": "[[LABEL_n]]",
                "warning": "This file is the re-identification KEY. It contains PHI. "
                           "Keep it on this machine. Never send it to an AI or upload it.",
                "mapping": mapping,
            },
            indent=2,
            ensure_ascii=False,
        ),
        encoding="utf-8",
    )
    return map_path


def process_file(src: Path, analyzer, out_dir: Path, explicit_out: Path | None,
                 used: set, reversible: bool = False) -> bool:
    text = extract(src)
    verify_extraction(text, src)  # fail closed before we write anything

    dest = explicit_out or _out_path(src, out_dir, used)
    dest.parent.mkdir(parents=True, exist_ok=True)

    if reversible:
        redacted, mapping = E.redact_reversible(text, analyzer)
        dest.write_text(redacted, encoding="utf-8")
        map_path = _write_map(dest, src, mapping)
        # count by label: "[[NAME_1]]" -> "NAME"
        counts = Counter(p[2:].rsplit("_", 1)[0] for p in mapping)
        summary = ", ".join(f"{k}:{v}" for k, v in sorted(counts.items())) or "nothing detected"
        print(f"  {src.name} -> {dest}")
        print(f"    {len(mapping)} unique identifiers tokenized ({summary})")
        print(f"    KEY -> {map_path}  [CONTAINS PHI: keep local, never send to an AI]")
    else:
        redacted, counts = redact_text(text, analyzer)
        dest.write_text(redacted, encoding="utf-8")
        summary = ", ".join(f"{k}:{v}" for k, v in sorted(counts.items())) or "nothing detected"
        print(f"  {src.name} -> {dest}")
        print(f"    {sum(counts.values())} spans redacted ({summary})")

    flags = residual_scan(redacted)
    if flags:
        print(
            f"    WARNING: {dest.name} may still contain {', '.join(flags)}-like text. "
            f"Line-wrapping in {src.suffix or 'the source'} extraction can split an "
            f"identifier so it slips past detection. Review this file by hand.",
            file=sys.stderr,
        )
    return bool(flags)


def gather_inputs(target: Path) -> list[Path]:
    if target.is_dir():
        files = sorted(p for p in target.iterdir() if p.suffix.lower() in SUPPORTED)
        if not files:
            raise SystemExit(f"No supported files ({', '.join(sorted(SUPPORTED))}) in {target}")
        return files
    if not target.exists():
        raise SystemExit(f"Not found: {target}")
    if target.suffix.lower() not in SUPPORTED:
        raise SystemExit(
            f"Unsupported type '{target.suffix}'. Supported: {', '.join(sorted(SUPPORTED))}"
        )
    return [target]


def require_installed() -> None:
    """Redaction runs offline, so it will not download anything on its own. If the
    one-time setup has not run, say so plainly instead of failing with a cryptic
    offline-cache error from deep inside the model loader."""
    if not MARKER.exists():
        raise SystemExit(
            "Setup has not run yet. This tool runs fully offline and will not "
            "download models on its own.\nRun the one-time installer first, from "
            "this skill's folder:\n    uv run scripts/bootstrap.py"
        )


def main() -> int:
    ap = argparse.ArgumentParser(description="De-identify PHI in a document.")
    ap.add_argument("input", help="file, folder, or '-' for stdin -> stdout")
    ap.add_argument(
        "--engine", choices=["baseline", "clinical"], default="clinical",
        help="clinical (default) adds the obi/deid_roberta_i2b2 transformer; "
             "baseline is spaCy + regex only",
    )
    ap.add_argument("-o", "--output", help="output file (single input only)")
    ap.add_argument("--out-dir", default="outputs", help="output folder (default: outputs/)")
    ap.add_argument(
        "--reversible", action="store_true",
        help="use unique placeholders ([[NAME_1]]) and write a <name>.map.json key so "
             "the text can be restored later with unredact.py. The key file contains PHI.",
    )
    args = ap.parse_args()

    require_installed()

    if args.reversible and args.input == "-":
        raise SystemExit("--reversible needs a file or folder (it writes a key file), not stdin.")

    # Status goes to stderr so stdout stays clean for the stdin -> stdout pipe.
    print(f"Loading the {args.engine} engine (first clinical run downloads the model once)...",
          file=sys.stderr)
    analyzer = E.build_analyzer(args.engine)

    # stdin -> stdout: for piping. No file is written.
    if args.input == "-":
        text = sys.stdin.read()
        verify_extraction(text, Path("<stdin>"))
        redacted, _ = redact_text(text, analyzer)
        sys.stdout.write(redacted)
        if residual_scan(redacted):
            print("WARNING: output may still contain identifier-like text; review it.",
                  file=sys.stderr)
        return 0

    target = Path(args.input)
    files = gather_inputs(target)

    if args.output and len(files) > 1:
        raise SystemExit("-o/--output works with a single input; use --out-dir for a folder.")

    out_dir = Path(args.out_dir)
    explicit_out = Path(args.output) if args.output else None

    mode = "reversible (with key files)" if args.reversible else "one-way"
    print(f"Redacting {len(files)} file(s), {mode}:")
    failures = warned = 0
    used: set = set()
    for src in files:
        try:
            if process_file(src, analyzer, out_dir, explicit_out, used, args.reversible):
                warned += 1
        except Exception as exc:  # fail closed per-file, keep going
            failures += 1
            print(f"  SKIPPED {src.name}: {exc}", file=sys.stderr)

    if warned:
        print(f"\n{warned} file(s) flagged with possible residual identifiers. Review them by hand.",
              file=sys.stderr)
    if failures:
        print(f"{failures} file(s) skipped without a redacted output. Review the messages above.",
              file=sys.stderr)
        return 1
    print("\nDone. Spot-check the output before trusting it.")
    if args.reversible:
        print("Reversible mode: send ONLY the .redacted.txt to the AI. Keep every "
              ".map.json key on this machine. Restore later with:  uv run unredact.py "
              "<ai_output> --map <name>.map.json")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
